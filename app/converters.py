import os
from pdf2image import convert_from_path
from pdf2docx import Converter
from docx import Document
from pptx import Presentation
import pandas as pd
from html import escape

def convert_pdf(input_path, conversion_type):
    """Perform the desired PDF conversion based on the conversion_type."""
    output_dir = 'app/outputs/'
    os.makedirs(output_dir, exist_ok=True)

    base_name = os.path.splitext(os.path.basename(input_path))[0]
    output_path = ""

    try:
        if conversion_type == 'pdf-doc':
            # Convert PDF to DOC
            output_path = os.path.join(output_dir, f"{base_name}.docx")
            cv = Converter(input_path)
            cv.convert(output_path)
            cv.close()

        elif conversion_type == 'pdf-xlsx':
            # Convert PDF table to XLSX
            output_path = os.path.join(output_dir, f"{base_name}.xlsx")
            data = [["Sample Data", "Row"], ["Another Row", "Value"]]  # Replace with actual extraction logic
            pd.DataFrame(data, columns=["Column1", "Column2"]).to_excel(output_path, index=False)

        elif conversion_type == 'pdf-pptx':
            # Convert PDF to PPTX
            output_path = os.path.join(output_dir, f"{base_name}.pptx")
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Title only layout
            slide.shapes.title.text = "Sample Slide from PDF"  # Replace with actual extraction logic
            presentation.save(output_path)

        elif conversion_type == 'pdf-html':
            # Convert PDF to HTML
            output_path = os.path.join(output_dir, f"{base_name}.html")
            html_content = f"<html><body><h1>Sample HTML</h1><p>Extracted from {escape(input_path)}</p></body></html>"
            with open(output_path, 'w') as html_file:
                html_file.write(html_content)

        elif conversion_type == 'pdf-tiff':
            # Convert PDF to TIFF
            images = convert_from_path(input_path)
            for i, img in enumerate(images):
                img.save(os.path.join(output_dir, f"{base_name}_{i + 1}.tiff"), 'TIFF')

        elif conversion_type == 'pdf-bmp':
            # Convert PDF to BMP
            images = convert_from_path(input_path)
            for i, img in enumerate(images):
                img.save(os.path.join(output_dir, f"{base_name}_{i + 1}.bmp"), 'BMP')

        elif conversion_type == 'pdf-rtf':
            # Convert PDF to RTF
            output_path = os.path.join(output_dir, f"{base_name}.rtf")
            document = Document()
            document.add_paragraph("Sample RTF content from PDF.")  # Replace with actual extraction logic
            document.save(output_path)

        elif conversion_type == 'pdf-jpg':
            # Convert PDF to JPG
            images = convert_from_path(input_path)
            for i, img in enumerate(images):
                img.save(os.path.join(output_dir, f"{base_name}_{i + 1}.jpg"), 'JPEG')

        elif conversion_type == 'pdf-png':
            # Convert PDF to PNG
            images = convert_from_path(input_path)
            for i, img in enumerate(images):
                img.save(os.path.join(output_dir, f"{base_name}_{i + 1}.png"), 'PNG')

        elif conversion_type == 'pdf-txt':
            # Extract text from PDF
            output_path = os.path.join(output_dir, f"{base_name}.txt")
            with open(output_path, 'w') as text_file:
                text_file.write("Sample extracted text.")  # Replace with actual extraction logic

        elif conversion_type == 'pdf-csv':
            # Convert PDF table to CSV
            output_path = os.path.join(output_dir, f"{base_name}.csv")
            data = [["Sample Data", "Row"], ["Another Row", "Value"]]  # Replace with actual extraction logic
            pd.DataFrame(data, columns=["Column1", "Column2"]).to_csv(output_path, index=False)

    except Exception as e:
        print(f"Error during conversion: {e}")
        return None

    return output_path or output_dir
